#!/usr/bin/env python3
"""
LuxeReserve - Advanced Database Presentation Generator (PPTX)
Generates a professional presentation for project defense.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# COLOR PALETTE
# ============================================================
DEEP_TEAL = RGBColor(0x14, 0x3D, 0x42)
WARM_BEIGE = RGBColor(0xF4, 0xED, 0xE3)
BRONZE = RGBColor(0xB5, 0x79, 0x3F)
TEXT_STRONG = RGBColor(0x17, 0x28, 0x2B)
TEXT_SOFT = RGBColor(0x60, 0x71, 0x72)
TEXT_ON_DARK = RGBColor(0xF7, 0xF1, 0xE7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)
ACCENT_BLUE = RGBColor(0x15, 0x65, 0xC0)

def add_bg(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    """Add a colored shape to the slide."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=TEXT_STRONG, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=14, color=TEXT_STRONG, spacing=Pt(6)):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = '\u2022  ' + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
    return txBox

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ============================================================
    # SLIDE 1: TITLE SLIDE
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide, DEEP_TEAL)
    
    # Decorative top bar
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), BRONZE)
    
    # Title
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11.333), Inches(1.5),
                 'LuxeReserve', font_size=54, bold=True, color=TEXT_ON_DARK, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(1), Inches(2.8), Inches(11.333), Inches(1),
                 'Advanced Database Technical Report', font_size=28, color=BRONZE, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11.333), Inches(0.8),
                 'Hotel Reservation Management System', font_size=20, color=TEXT_ON_DARK, alignment=PP_ALIGN.CENTER)
    
    # Divider line
    add_shape(slide, Inches(4), Inches(4.8), Inches(5.333), Inches(0.03), BRONZE)
    
    add_text_box(slide, Inches(1), Inches(5.2), Inches(11.333), Inches(0.6),
                 'Polyglot Persistence: SQL Server 2022 + MongoDB Atlas', font_size=16, color=TEXT_ON_DARK, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(1), Inches(5.8), Inches(11.333), Inches(0.5),
                 'Course: Advanced Database  |  Group: DAF04  |  April 2026', font_size=14, color=TEXT_SOFT, alignment=PP_ALIGN.CENTER)
    
    # Decorative bottom bar
    add_shape(slide, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), BRONZE)
    
    # ============================================================
    # SLIDE 2: AGENDA
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DEEP_TEAL)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
                 'Presentation Agenda', font_size=36, bold=True, color=TEXT_ON_DARK)
    
    items = [
        'Project Introduction & Objectives',
        'Technology Stack Overview',
        'System Architecture (Three-Tier + Hybrid Data Layer)',
        'Database Design (30 SQL Tables + 3 MongoDB Collections)',
        'Advanced Database Techniques (CTE, Locking, Window Functions, Triggers, SP, Views)',
        'System Features & Functionality',
        'Source Code Architecture',
        'Conclusion & Future Work',
    ]
    add_bullet_text(slide, Inches(1.5), Inches(1.8), Inches(10.333), Inches(5), items, font_size=18, spacing=Pt(10))
    
    # ============================================================
    # SLIDE 3: PROJECT OVERVIEW
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DEEP_TEAL)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
                 '1. Project Introduction', font_size=36, bold=True, color=TEXT_ON_DARK)
    
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.733), Inches(1.2),
                 'LuxeReserve is a comprehensive Hotel Reservation Management System (HRMS) designed for '
                 'luxury hotel chains operating across Southeast Asia. It provides end-to-end solutions for '
                 'managing hotel inventory, guest reservations, payment processing, front desk operations, '
                 'housekeeping, maintenance, and revenue analytics.',
                 font_size=16, color=TEXT_STRONG)
    
    # Key metrics boxes
    metrics = [
        ('30', 'SQL Tables'),
        ('3', 'MongoDB Collections'),
        ('82', 'API Endpoints'),
        ('12', 'Page Routes'),
        ('14', 'Route Groups'),
        ('3', 'User Types'),
    ]
    for i, (num, label) in enumerate(metrics):
        col = i % 3
        row = i // 3
        left = Inches(0.8 + col * 4.2)
        top = Inches(3.5 + row * 1.8)
        
        box = add_shape(slide, left, top, Inches(3.8), Inches(1.5), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text_box(slide, left + Inches(0.2), top + Inches(0.1), Inches(3.4), Inches(0.8),
                     num, font_size=40, bold=True, color=DEEP_TEAL, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.2), top + Inches(0.9), Inches(3.4), Inches(0.5),
                     label, font_size=16, color=TEXT_SOFT, alignment=PP_ALIGN.CENTER)
    
    # ============================================================
    # SLIDE 4: TECHNOLOGY STACK
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DEEP_TEAL)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
                 '2. Technology Stack', font_size=36, bold=True, color=TEXT_ON_DARK)
    
    # Frontend
    add_shape(slide, Inches(0.5), Inches(1.6), Inches(3.8), Inches(5.5), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(0.7), Inches(1.7), Inches(3.4), Inches(0.5),
                 'Frontend', font_size=22, bold=True, color=DEEP_TEAL, alignment=PP_ALIGN.CENTER)
    fe_items = [
        'React 18 + Vite',
        'React Router v6',
        'Context API (Auth, Flash)',
        '12 Page Components',
        'Premium UI Design',
        'Warm Beige + Deep Teal',
    ]
    add_bullet_text(slide, Inches(0.7), Inches(2.4), Inches(3.4), Inches(4.5), fe_items, font_size=14, spacing=Pt(8))
    
    # Backend
    add_shape(slide, Inches(4.7), Inches(1.6), Inches(3.8), Inches(5.5), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(4.9), Inches(1.7), Inches(3.4), Inches(0.5),
                 'Backend', font_size=22, bold=True, color=DEEP_TEAL, alignment=PP_ALIGN.CENTER)
    be_items = [
        'Node.js + Express',
        '82 RESTful Endpoints',
        'JWT Authentication',
        'bcryptjs Password Hashing',
        'VNPay Integration',
        'Nodemailer (Gmail SMTP)',
    ]
    add_bullet_text(slide, Inches(4.9), Inches(2.4), Inches(3.4), Inches(4.5), be_items, font_size=14, spacing=Pt(8))
    
    # Database
    add_shape(slide, Inches(8.9), Inches(1.6), Inches(3.8), Inches(5.5), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(9.1), Inches(1.7), Inches(3.4), Inches(0.5),
                 'Database', font_size=22, bold=True, color=DEEP_TEAL, alignment=PP_ALIGN.CENTER)
    db_items = [
        'SQL Server 2022 Express',
        'MongoDB Atlas',
        'Polyglot Persistence',
        '30 Tables / 3 Collections',
        'ACID Transactions',
        'Flexible Schema Content',
    ]
    add_bullet_text(slide, Inches(9.1), Inches(2.4), Inches(3.4), Inches(4.5), db_items, font_size=14, spacing=Pt(8))
    
    # ============================================================
    # SLIDE 5: SYSTEM ARCHITECTURE
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DEEP_TEAL)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
                 '3. System Architecture', font_size=36, bold=True, color=TEXT_ON_DARK)
    
    # Three-Tier Architecture Diagram
    tiers = [
        ('Presentation Tier', 'React SPA\nPort 5173', DEEP_TEAL),
        ('Application Tier', 'Node.js Express\nPort 3000', BRONZE),
        ('Data Tier', 'SQL Server + MongoDB', RGBColor(0x2E, 0x7D, 0x32)),
    ]
    for i, (title, desc, color) in enumerate(tiers):
        left = Inches(0.8 + i * 4.2)
        box = add_shape(slide, left, Inches(1.6), Inches(3.8), Inches(2.5), color, MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text_box(slide, left + Inches(0.2), Inches(1.7), Inches(3.4), Inches(0.5),
                     title, font_size=20, bold=True, color=TEXT_ON_DARK, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.2), Inches(2.4), Inches(3.4), Inches(1.5),
                     desc, font_size=16, color=TEXT_ON_DARK, alignment=PP_ALIGN.CENTER)
        
        # Arrow between tiers
        if i < 2:
            add_text_box(slide, left + Inches(3.8), Inches(2.4), Inches(0.4), Inches(0.5),
                         '\u25B6', font_size=24, color=BRONZE, alignment=PP_ALIGN.CENTER)
    
    # Hybrid Data Layer
    add_shape(slide, Inches(0.5), Inches(4.5), Inches(12.333), Inches(2.7), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(0.7), Inches(4.6), Inches(11.933), Inches(0.5),
                 'Hybrid Data Layer Design', font_size=20, bold=True, color=DEEP_TEAL)
    
    hybrid_items = [
        'SQL Server handles all transactional data: reservations, payments, inventory, guest profiles, audit logs',
        'MongoDB stores rich content: hotel descriptions, image galleries, amenity features, room type details',
        'API layer merges data from both databases using business key mapping (hotel_id, amenity_code, room_type_code)',
        '94% of endpoints served primarily from SQL Server; 6% require hybrid data merging',
    ]
    add_bullet_text(slide, Inches(0.7), Inches(5.2), Inches(11.933), Inches(3.5), hybrid_items, font_size=14, spacing=Pt(6))
    
    # ============================================================
    # SLIDE 6: DATABASE DESIGN
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DEEP_TEAL)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
                 '4. Database Design', font_size=36, bold=True, color=TEXT_ON_DARK)
    
    # 12 Domains
    domains = [
        'Location Hierarchy', 'Hotel Organization', 'Room Management',
        'Guest Management', 'System Users', 'Pricing & Promotions',
        'Reservation Core', 'Financial', 'Stay & Services',
        'Operations', 'Loyalty', 'Audit & Locks',
    ]
    for i, domain in enumerate(domains):
        col = i % 4
        row = i // 4
        left = Inches(0.5 + col * 3.2)
        top = Inches(1.5 + row * 1.2)
        
        box = add_shape(slide, left, top, Inches(2.9), Inches(0.9), DEEP_TEAL, MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text_box(slide, left + Inches(0.1), top + Inches(0.15), Inches(2.7), Inches(0.6),
                     domain, font_size=14, bold=True, color=TEXT_ON_DARK, alignment=PP_ALIGN.CENTER)
    
    # Key design principles
    add_shape(slide, Inches(0.5), Inches(5.2), Inches(12.333), Inches(2.0), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(0.7), Inches(5.3), Inches(11.933), Inches(0.4),
                 'Key Design Principles', font_size=18, bold=True, color=DEEP_TEAL)
    principles = [
        'Computed Persisted Column: Guest.full_name (CONCAT + COALESCE)',
        'Denormalized hotel_id in RoomAvailability and HousekeepingTask for query performance',
        'Snapshot currency_code in Reservation, Payment, Invoice for historical accuracy',
        'CHECK constraint on RoomFeature: at least one FK (room_id or room_type_id) must be NOT NULL',
    ]
    add_bullet_text(slide, Inches(0.7), Inches(5.8), Inches(11.933), Inches(1.5), principles, font_size=13, spacing=Pt(4))
    
    # ============================================================
    # SLIDE 7: ADVANCED TECHNIQUES - CTE & LOCKING
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DEEP_TEAL)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
                 '5. Advanced Database Techniques (Part 1)', font_size=36, bold=True, color=TEXT_ON_DARK)
    
    # Recursive CTE
    add_shape(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.8), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(0.4),
                 'Recursive CTE - Location Hierarchy', font_size=18, bold=True, color=DEEP_TEAL)
    cte_items = [
        'Adjacency list model: 5 levels (Region, Country, State, City, District)',
        'Anchor member selects root; recursive member joins children',
        'vw_LocationTree view with depth tracking and hotel_count aggregation',
        'Enables queries like "find all hotels in Southeast Asia" in one query',
    ]
    add_bullet_text(slide, Inches(0.7), Inches(2.1), Inches(5.6), Inches(2.0), cte_items, font_size=13, spacing=Pt(4))
    
    # Pessimistic Locking
    add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.8), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.4),
                 'Pessimistic Locking - Double-Booking Prevention', font_size=18, bold=True, color=DEEP_TEAL)
    pl_items = [
        'UPDLOCK + HOLDLOCK hints on RoomAvailability',
        'sp_ReserveRoom: lock, check, update, commit',
        'Second transaction BLOCKED until first completes',
        'sp_TransferRoom: atomic release old + lock new room',
    ]
    add_bullet_text(slide, Inches(7.0), Inches(2.1), Inches(5.6), Inches(2.0), pl_items, font_size=13, spacing=Pt(4))
    
    # Optimistic Locking
    add_shape(slide, Inches(0.5), Inches(4.6), Inches(6.0), Inches(2.6), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(0.7), Inches(4.7), Inches(5.6), Inches(0.4),
                 'Optimistic Locking - Inventory Updates', font_size=18, bold=True, color=DEEP_TEAL)
    ol_items = [
        'version_no column for conflict detection',
        'UPDATE ... WHERE version_no = @expected_version',
        'Returns HTTP 409 Conflict on version mismatch',
        'Suitable for low-contention admin operations',
    ]
    add_bullet_text(slide, Inches(0.7), Inches(5.2), Inches(5.6), Inches(1.8), ol_items, font_size=13, spacing=Pt(4))
    
    # Window Functions
    add_shape(slide, Inches(6.8), Inches(4.6), Inches(6.0), Inches(2.6), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(7.0), Inches(4.7), Inches(5.6), Inches(0.4),
                 'Window Functions - Revenue Analytics', font_size=18, bold=True, color=DEEP_TEAL)
    wf_items = [
        'SUM() OVER(PARTITION BY) for monthly revenue per hotel',
        'DENSE_RANK() for top-performing room types',
        'Cumulative running totals with ROWS BETWEEN',
        'Revenue share percentage calculations',
    ]
    add_bullet_text(slide, Inches(7.0), Inches(5.2), Inches(5.6), Inches(1.8), wf_items, font_size=13, spacing=Pt(4))
    
    # ============================================================
    # SLIDE 8: ADVANCED TECHNIQUES - TRIGGERS, SP, VIEWS
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DEEP_TEAL)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
                 '5. Advanced Database Techniques (Part 2)', font_size=36, bold=True, color=TEXT_ON_DARK)
    
    # Triggers
    add_shape(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.8), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(0.4),
                 'Triggers - Automated Audit & Integrity', font_size=18, bold=True, color=DEEP_TEAL)
    tr_items = [
        'trg_Payment_AuditLog: logs payment creation and status changes',
        'trg_Guest_AuditLog: captures profile changes (email, name, VIP)',
        'trg_GuestAuth_AuditLog: detects password changes, email verification',
        'trg_RoomRate_PriceIntegrityGuard: flags rate changes > 50% as CRITICAL',
        'All triggers use TRY-CATCH with explicit transaction management',
    ]
    add_bullet_text(slide, Inches(0.7), Inches(2.1), Inches(5.6), Inches(2.0), tr_items, font_size=13, spacing=Pt(4))
    
    # Stored Procedures
    add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.8), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.4),
                 'Stored Procedures - Business Logic', font_size=18, bold=True, color=DEEP_TEAL)
    sp_items = [
        'sp_ReserveRoom: pessimistic locking for booking',
        'sp_TransferRoom: atomic room transfer with lock management',
        'sp_CreateReservation: orchestrates full booking creation',
        'SET XACT_ABORT ON for automatic rollback',
        'Comprehensive TRY-CATCH error handling',
    ]
    add_bullet_text(slide, Inches(7.0), Inches(2.1), Inches(5.6), Inches(2.0), sp_items, font_size=13, spacing=Pt(4))
    
    # Views
    add_shape(slide, Inches(0.5), Inches(4.6), Inches(6.0), Inches(2.6), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(0.7), Inches(4.7), Inches(5.6), Inches(0.4),
                 'Views - Financial & Analytical', font_size=18, bold=True, color=DEEP_TEAL)
    vw_items = [
        'vw_ReservationTotal: single source of truth for financials',
        'Aggregates from ReservationRoom, ReservationService, Payment',
        'vw_LocationTree: hierarchical location queries',
        'vw_RevenueByHotel: monthly revenue with window functions',
    ]
    add_bullet_text(slide, Inches(0.7), Inches(5.2), Inches(5.6), Inches(1.8), vw_items, font_size=13, spacing=Pt(4))
    
    # Computed Columns
    add_shape(slide, Inches(6.8), Inches(4.6), Inches(6.0), Inches(2.6), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(7.0), Inches(4.7), Inches(5.6), Inches(0.4),
                 'Computed Columns - Derived Data', font_size=18, bold=True, color=DEEP_TEAL)
    cc_items = [
        'Guest.full_name: CONCAT + COALESCE with PERSISTED',
        'Physically stored for indexability',
        'Eliminates application-level name concatenation',
        'Ensures consistent name formatting across all interfaces',
    ]
    add_bullet_text(slide, Inches(7.0), Inches(5.2), Inches(5.6), Inches(1.8), cc_items, font_size=13, spacing=Pt(4))
    
    # ============================================================
    # SLIDE 9: SYSTEM FEATURES
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DEEP_TEAL)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
                 '6. System Features', font_size=36, bold=True, color=TEXT_ON_DARK)
    
    features = [
        ('Reservation Management', 'ACID booking lifecycle: PENDING \u2192 CONFIRMED \u2192 CHECKED_IN \u2192 CHECKED_OUT. Room transfers, cancellations, status history tracking.'),
        ('Payment Processing', '5 payment types (Deposit, Prepayment, Full, Refund, Incidental). VNPay integration. Validation against vw_ReservationTotal.'),
        ('Front Desk Operations', 'Check-in/check-out workflows. StayRecord creation. Auto housekeeping task generation. Guest lookup and room transfers.'),
        ('Housekeeping & Maintenance', 'Priority-based cleaning tasks. Staff assignment. Severity-based maintenance tickets. Auto room status updates.'),
        ('Guest Services & Loyalty', 'Service ordering (spa, dining, transfer). Loyalty tiers (Silver to Black). Points redemption for promotional vouchers.'),
        ('Admin Dashboard', 'Revenue analytics with window functions. Price Guard monitoring. Invoice generation. User management and role assignment.'),
    ]
    for i, (title, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        left = Inches(0.5 + col * 4.2)
        top = Inches(1.5 + row * 2.9)
        
        add_shape(slide, left, top, Inches(3.9), Inches(2.6), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text_box(slide, left + Inches(0.2), top + Inches(0.1), Inches(3.5), Inches(0.4),
                     title, font_size=16, bold=True, color=DEEP_TEAL)
        add_text_box(slide, left + Inches(0.2), top + Inches(0.6), Inches(3.5), Inches(1.8),
                     desc, font_size=12, color=TEXT_STRONG)
    
    # ============================================================
    # SLIDE 10: SOURCE CODE ARCHITECTURE
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DEEP_TEAL)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
                 '7. Source Code Architecture', font_size=36, bold=True, color=TEXT_ON_DARK)
    
    # Backend
    add_shape(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.8), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(0.4),
                 'Backend (Node.js/Express)', font_size=18, bold=True, color=DEEP_TEAL)
    be_items = [
        'src/app.js: Main entry point, middleware config, route mounting',
        'src/config/: SQL Server + MongoDB connection configurations',
        'src/middleware/: JWT auth stack (attachAuthContext, requireAuth)',
        'src/routes/: 14 route modules (hotels, rooms, guests, auth, etc.)',
        'src/services/: mail.js (Nodemailer), vnpay.js (payment gateway)',
        'database/sql/: 24 numbered migration scripts (01-24)',
    ]
    add_bullet_text(slide, Inches(0.7), Inches(2.1), Inches(5.6), Inches(2.0), be_items, font_size=12, spacing=Pt(3))
    
    # Frontend
    add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.8), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.4),
                 'Frontend (React/Vite)', font_size=18, bold=True, color=DEEP_TEAL)
    fe_items = [
        'frontend/src/main.jsx: Entry point with context providers',
        'frontend/src/App.jsx: React Router v6 route definitions',
        'frontend/src/pages/: 12 page components (Dashboard, Search, etc.)',
        'frontend/src/context/: AuthContext, FlashContext for state management',
        'frontend/src/lib/api.js: HTTP client with JWT token attachment',
        'frontend/src/utils/: Formatters, hotel image management',
    ]
    add_bullet_text(slide, Inches(7.0), Inches(2.1), Inches(5.6), Inches(2.0), fe_items, font_size=12, spacing=Pt(3))
    
    # Testing
    add_shape(slide, Inches(0.5), Inches(4.6), Inches(12.333), Inches(2.6), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(0.7), Inches(4.7), Inches(11.933), Inches(0.4),
                 'Testing & Quality Assurance', font_size=18, bold=True, color=DEEP_TEAL)
    test_items = [
        'Playwright end-to-end testing for critical user workflows (booking, payment, check-in/check-out)',
        'Concurrent booking test (test_concurrent_booking.js): validates pessimistic locking prevents double-booking',
        'Reservation payment flow test (test_reservation_payment_flow.js): validates complete payment lifecycle',
        'check_remaining.py: automated validation of project requirements against rubric',
        'Playwright HTML report in playwright-report/ for visual test result review',
    ]
    add_bullet_text(slide, Inches(0.7), Inches(5.2), Inches(11.933), Inches(1.8), test_items, font_size=13, spacing=Pt(4))
    
    # ============================================================
    # SLIDE 11: CONCLUSION
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DEEP_TEAL)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.8),
                 '8. Conclusion', font_size=36, bold=True, color=TEXT_ON_DARK)
    
    # Achievements
    add_shape(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.8), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(0.4),
                 'Key Achievements', font_size=18, bold=True, color=DEEP_TEAL)
    ach_items = [
        'ACID-compliant transactions with pessimistic locking preventing double-booking',
        'Polyglot persistence: SQL Server (30 tables) + MongoDB (3 collections)',
        '82 RESTful API endpoints across 14 route groups',
        'React-based UI for 3 user groups (guests, staff, admins)',
        'Automated audit logging through 4 database triggers',
        'Complete database setup pipeline with 24 migration scripts',
    ]
    add_bullet_text(slide, Inches(0.7), Inches(2.1), Inches(5.6), Inches(2.0), ach_items, font_size=12, spacing=Pt(3))
    
    # Challenges
    add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.8), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.4),
                 'Challenges & Lessons Learned', font_size=18, bold=True, color=DEEP_TEAL)
    ch_items = [
        'Implementing pessimistic locking without deadlocks in high-concurrency scenarios',
        'Designing hybrid data layer with efficient SQL + MongoDB merging',
        'Ensuring data consistency across two database systems',
        'Comprehensive error handling in triggers and stored procedures',
        'Managing transaction boundaries across multiple database operations',
    ]
    add_bullet_text(slide, Inches(7.0), Inches(2.1), Inches(5.6), Inches(2.0), ch_items, font_size=12, spacing=Pt(3))
    
    # Future Work
    add_shape(slide, Inches(0.5), Inches(4.6), Inches(12.333), Inches(2.6), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(0.7), Inches(4.7), Inches(11.933), Inches(0.4),
                 'Future Enhancements', font_size=18, bold=True, color=DEEP_TEAL)
    fw_items = [
        'Integration with third-party OTAs (Booking.com, Agoda) for centralized inventory management',
        'Native mobile applications for guests and staff with push notifications',
        'AI-based revenue management for dynamic room pricing optimization',
        'Multi-property consolidated accounting for enterprise-level financial reporting',
        'Integration with Property Management Systems (PMS) and Point-of-Sale (POS) systems',
    ]
    add_bullet_text(slide, Inches(0.7), Inches(5.2), Inches(11.933), Inches(1.8), fw_items, font_size=13, spacing=Pt(4))
    
    # ============================================================
    # SLIDE 12: THANK YOU
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DEEP_TEAL)
    
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), BRONZE)
    add_shape(slide, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), BRONZE)
    
    add_text_box(slide, Inches(1), Inches(2.0), Inches(11.333), Inches(1.5),
                 'Thank You', font_size=54, bold=True, color=TEXT_ON_DARK, alignment=PP_ALIGN.CENTER)
    
    add_shape(slide, Inches(4), Inches(3.5), Inches(5.333), Inches(0.03), BRONZE)
    
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11.333), Inches(1),
                 'Questions & Discussion', font_size=28, color=BRONZE, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(1), Inches(5.0), Inches(11.333), Inches(0.8),
                 'LuxeReserve - Advanced Database Technical Report\nGroup DAF04 | April 2026',
                 font_size=16, color=TEXT_ON_DARK, alignment=PP_ALIGN.CENTER)
    
    # ============================================================
    # SAVE
    # ============================================================
    output_path = 'LuxeReserve_Advanced_Database_Presentation.pptx'
    prs.save(output_path)
    print('[OK] Presentation saved to: %s' % output_path)
    return output_path

if __name__ == '__main__':
    create_presentation()
